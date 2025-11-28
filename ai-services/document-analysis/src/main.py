"""
文档分析AI服务
提供失效分析相关的文档处理和分析功能
"""

from fastapi import FastAPI, HTTPException, UploadFile, File, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Dict, Any, Optional, Union
import uvicorn
import asyncio
import redis
import json
import io
import base64
from datetime import datetime
from pathlib import Path
import tempfile
import os
from loguru import logger

# 文档处理库
try:
    import PyPDF2
    import pdfplumber
    from docx import Document
    from pptx import Presentation
    import openpyxl
    import pytesseract
    from PIL import Image
    import easyocr
except ImportError as e:
    logger.warning(f"某些文档处理库未安装: {e}")

# 文本处理
import nltk
import jieba
import re
from collections import Counter

# 配置日志
logger.add("logs/document_analysis.log", rotation="10 MB", level="INFO")

app = FastAPI(
    title="文档分析AI服务",
    description="失效分析智能辅助平台 - 文档分析服务",
    version="1.0.0"
)

# CORS配置
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Redis连接
redis_client = redis.Redis(host='redis', port=6379, db=0, decode_responses=True)

# 初始化OCR
try:
    ocr_reader = easyocr.Reader(['ch_sim', 'en'])
    logger.info("EasyOCR 初始化成功")
except Exception as e:
    logger.warning(f"EasyOCR 初始化失败: {e}")
    ocr_reader = None

# 数据模型
class DocumentAnalysisRequest(BaseModel):
    task_id: str
    analysis_type: str
    parameters: Dict[str, Any] = {}

class DocumentAnalysisResult(BaseModel):
    task_id: str
    status: str
    results: Dict[str, Any]
    execution_time: float
    timestamp: str

class TextExtractionRequest(BaseModel):
    extract_type: str  # "text", "tables", "images", "metadata"
    language: str = "auto"
    ocr_enabled: bool = True

class DocumentSummaryRequest(BaseModel):
    summary_type: str  # "brief", "detailed", "key_points"
    max_length: Optional[int] = None
    language: str = "zh"

# 文档分析核心功能
class DocumentAnalyzer:
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.doc', '.pptx', '.ppt', '.xlsx', '.xls', '.txt', '.rtf']
        self.temp_dir = Path(tempfile.gettempdir()) / "smartfa_docs"
        self.temp_dir.mkdir(exist_ok=True)
    
    async def analyze_document(self, file_data: bytes, filename: str, analysis_type: str, parameters: Dict) -> Dict[str, Any]:
        """分析文档"""
        try:
            # 保存临时文件
            temp_path = await self.save_temp_file(file_data, filename)
            file_ext = Path(filename).suffix.lower()
            
            results = {}
            
            if analysis_type == "text_extraction":
                results = await self.extract_text(temp_path, file_ext, parameters)
            elif analysis_type == "document_summary":
                results = await self.generate_summary(temp_path, file_ext, parameters)
            elif analysis_type == "keyword_extraction":
                results = await self.extract_keywords(temp_path, file_ext, parameters)
            elif analysis_type == "document_classification":
                results = await self.classify_document(temp_path, file_ext, parameters)
            elif analysis_type == "entity_extraction":
                results = await self.extract_entities(temp_path, file_ext, parameters)
            elif analysis_type == "document_similarity":
                results = await self.calculate_similarity(temp_path, file_ext, parameters)
            else:
                raise ValueError(f"不支持的分析类型: {analysis_type}")
            
            # 清理临时文件
            await self.cleanup_temp_file(temp_path)
            
            return results
            
        except Exception as e:
            logger.error(f"文档分析失败: {str(e)}")
            raise
    
    async def save_temp_file(self, file_data: bytes, filename: str) -> Path:
        """保存临时文件"""
        temp_path = self.temp_dir / f"{datetime.now().timestamp()}_{filename}"
        with open(temp_path, 'wb') as f:
            f.write(file_data)
        return temp_path
    
    async def cleanup_temp_file(self, temp_path: Path):
        """清理临时文件"""
        try:
            if temp_path.exists():
                temp_path.unlink()
        except Exception as e:
            logger.warning(f"清理临时文件失败: {e}")
    
    async def extract_text(self, file_path: Path, file_ext: str, parameters: Dict) -> Dict[str, Any]:
        """提取文档文本"""
        text_content = ""
        tables = []
        images = []
        metadata = {}
        
        extract_type = parameters.get("extract_type", "text")
        ocr_enabled = parameters.get("ocr_enabled", True)
        language = parameters.get("language", "auto")
        
        try:
            if file_ext == '.pdf':
                text_content, tables, images, metadata = await self.extract_from_pdf(file_path, ocr_enabled)
            elif file_ext in ['.docx', '.doc']:
                text_content, tables, images, metadata = await self.extract_from_docx(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                text_content, images, metadata = await self.extract_from_pptx(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                text_content, tables, metadata = await self.extract_from_excel(file_path)
            elif file_ext == '.txt':
                text_content = await self.extract_from_txt(file_path)
            else:
                raise ValueError(f"不支持的文件格式: {file_ext}")
            
            # 文本清理和预处理
            if text_content:
                text_content = await self.clean_text(text_content)
                if language == "zh":
                    text_content = await self.segment_chinese_text(text_content)
            
            results = {
                "text_length": len(text_content),
                "word_count": len(text_content.split()) if text_content else 0,
                "language_detected": await self.detect_language(text_content) if text_content else "unknown"
            }
            
            if extract_type in ["text", "all"]:
                results["text_content"] = text_content[:10000] + "..." if len(text_content) > 10000 else text_content
            
            if extract_type in ["tables", "all"]:
                results["tables"] = tables[:10]  # 限制返回表格数量
            
            if extract_type in ["images", "all"]:
                results["images"] = images[:20]  # 限制返回图片数量
            
            if extract_type in ["metadata", "all"]:
                results["metadata"] = metadata
            
            return results
            
        except Exception as e:
            logger.error(f"文本提取失败: {str(e)}")
            raise
    
    async def extract_from_pdf(self, file_path: Path, ocr_enabled: bool) -> tuple:
        """从PDF提取内容"""
        text_content = ""
        tables = []
        images = []
        metadata = {}
        
        try:
            # 使用pdfplumber提取文本和表格
            with pdfplumber.open(file_path) as pdf:
                metadata = {
                    "title": pdf.metadata.get('Title', ''),
                    "author": pdf.metadata.get('Author', ''),
                    "creator": pdf.metadata.get('Creator', ''),
                    "producer": pdf.metadata.get('Producer', ''),
                    "creation_date": str(pdf.metadata.get('CreationDate', '')),
                    "modification_date": str(pdf.metadata.get('ModDate', '')),
                    "page_count": len(pdf.pages)
                }
                
                for page_num, page in enumerate(pdf.pages):
                    # 提取文本
                    page_text = page.extract_text()
                    if page_text:
                        text_content += f"\n--- 第{page_num + 1}页 ---\n{page_text}\n"
                    
                    # 提取表格
                    page_tables = page.extract_tables()
                    for table in page_tables:
                        if table:
                            tables.append({
                                "page": page_num + 1,
                                "rows": len(table),
                                "cols": len(table[0]) if table else 0,
                                "data": table
                            })
                    
                    # 提取图片
                    if hasattr(page, 'images'):
                        for img in page.images:
                            images.append({
                                "page": page_num + 1,
                                "x0": img.get('x0', 0),
                                "y0": img.get('y0', 0),
                                "x1": img.get('x1', 0),
                                "y1": img.get('y1', 0),
                                "width": img.get('width', 0),
                                "height": img.get('height', 0)
                            })
            
            # 如果OCR启用且文本内容较少，尝试OCR
            if ocr_enabled and len(text_content.strip()) < 100 and ocr_reader:
                ocr_text = await self.ocr_pdf(file_path)
                if ocr_text:
                    text_content += f"\n--- OCR识别内容 ---\n{ocr_text}"
        
        except Exception as e:
            logger.error(f"PDF提取失败: {str(e)}")
            raise
        
        return text_content, tables, images, metadata
    
    async def extract_from_docx(self, file_path: Path) -> tuple:
        """从DOCX提取内容"""
        text_content = ""
        tables = []
        images = []
        metadata = {}
        
        try:
            doc = Document(file_path)
            
            # 提取文档属性
            core_props = doc.core_properties
            metadata = {
                "title": core_props.title or '',
                "author": core_props.author or '',
                "subject": core_props.subject or '',
                "keywords": core_props.keywords or '',
                "comments": core_props.comments or '',
                "category": core_props.category or '',
                "created": str(core_props.created) if core_props.created else '',
                "modified": str(core_props.modified) if core_props.modified else '',
                "revision": str(core_props.revision) if core_props.revision else ''
            }
            
            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"
            
            # 提取表格
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                
                tables.append({
                    "table_index": table_idx + 1,
                    "rows": len(table_data),
                    "cols": len(table_data[0]) if table_data else 0,
                    "data": table_data
                })
            
            # 统计图片数量（DOCX中的图片提取较复杂，这里只统计）
            image_count = 0
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_count += 1
            
            if image_count > 0:
                images.append({"count": image_count, "note": "图片提取需要额外处理"})
        
        except Exception as e:
            logger.error(f"DOCX提取失败: {str(e)}")
            raise
        
        return text_content, tables, images, metadata
    
    async def extract_from_pptx(self, file_path: Path) -> tuple:
        """从PPTX提取内容"""
        text_content = ""
        images = []
        metadata = {}
        
        try:
            prs = Presentation(file_path)
            
            # 提取文档属性
            core_props = prs.core_properties
            metadata = {
                "title": core_props.title or '',
                "author": core_props.author or '',
                "subject": core_props.subject or '',
                "keywords": core_props.keywords or '',
                "comments": core_props.comments or '',
                "created": str(core_props.created) if core_props.created else '',
                "modified": str(core_props.modified) if core_props.modified else '',
                "slide_count": len(prs.slides)
            }
            
            # 提取每页内容
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = f"\n--- 幻灯片 {slide_idx + 1} ---\n"
                
                # 提取文本框内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                    
                    # 统计图片
                    if shape.shape_type == 13:  # 图片类型
                        images.append({
                            "slide": slide_idx + 1,
                            "type": "image",
                            "name": getattr(shape, 'name', f'image_{len(images)}')
                        })
                
                text_content += slide_text
        
        except Exception as e:
            logger.error(f"PPTX提取失败: {str(e)}")
            raise
        
        return text_content, images, metadata
    
    async def extract_from_excel(self, file_path: Path) -> tuple:
        """从Excel提取内容"""
        text_content = ""
        tables = []
        metadata = {}
        
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            metadata = {
                "sheet_count": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames,
                "active_sheet": workbook.active.title if workbook.active else ''
            }
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_data = []
                sheet_text = f"\n--- 工作表: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    sheet_data.append(row_data)
                    
                    # 提取有文本的单元格
                    row_text = ' | '.join([cell for cell in row_data if cell.strip()])
                    if row_text.strip():
                        sheet_text += row_text + "\n"
                
                text_content += sheet_text
                
                if sheet_data:
                    tables.append({
                        "sheet_name": sheet_name,
                        "rows": len(sheet_data),
                        "cols": len(sheet_data[0]) if sheet_data else 0,
                        "data": sheet_data
                    })
        
        except Exception as e:
            logger.error(f"Excel提取失败: {str(e)}")
            raise
        
        return text_content, tables, metadata
    
    async def extract_from_txt(self, file_path: Path) -> str:
        """从TXT提取内容"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试其他编码
            for encoding in ['gbk', 'gb2312', 'latin-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            raise ValueError("无法解码文本文件")
    
    async def clean_text(self, text: str) -> str:
        """清理文本"""
        if not text:
            return ""
        
        # 移除多余空白
        text = re.sub(r'\s+', ' ', text)
        
        # 移除特殊字符（保留中文、英文、数字、基本标点）
        text = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9\s.,!?;:()[\]{}"\'-]', '', text)
        
        # 移除多余换行
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        return text.strip()
    
    async def segment_chinese_text(self, text: str) -> str:
        """中文分词"""
        try:
            # 使用jieba分词
            words = jieba.lcut(text)
            return ' '.join(words)
        except Exception as e:
            logger.warning(f"中文分词失败: {e}")
            return text
    
    async def detect_language(self, text: str) -> str:
        """检测语言"""
        try:
            from langdetect import detect
            return detect(text)
        except Exception as e:
            logger.warning(f"语言检测失败: {e}")
            return "unknown"
    
    async def ocr_pdf(self, file_path: Path) -> str:
        """PDF OCR识别"""
        if not ocr_reader:
            return ""
        
        try:
            # 将PDF转换为图片进行OCR
            # 这里简化处理，实际应该使用pdf2image等库
            result = ocr_reader.readtext(str(file_path))
            text = " ".join([item[1] for item in result])
            return text
        except Exception as e:
            logger.error(f"PDF OCR失败: {str(e)}")
            return ""
    
    async def generate_summary(self, file_path: Path, file_ext: str, parameters: Dict) -> Dict[str, Any]:
        """生成文档摘要"""
        # 首先提取文本
        text_content, _, _, metadata = await self.extract_text(file_path, file_ext, {"extract_type": "text"})
        
        if not text_content:
            return {"error": "无法提取文档内容"}
        
        summary_type = parameters.get("summary_type", "brief")
        max_length = parameters.get("max_length", 200)
        language = parameters.get("language", "zh")
        
        # 简单的摘要生成（实际应用中应使用更高级的NLP模型）
        sentences = text_content.split('。')
        sentences = [s.strip() for s in sentences if s.strip()]
        
        if summary_type == "brief":
            # 简短摘要：取前几句
            summary_sentences = sentences[:3]
        elif summary_type == "detailed":
            # 详细摘要：取更多句子
            summary_sentences = sentences[:10]
        elif summary_type == "key_points":
            # 关键点：选择较长的句子
            summary_sentences = sorted(sentences, key=len, reverse=True)[:5]
        else:
            summary_sentences = sentences[:5]
        
        summary = '。'.join(summary_sentences)
        
        # 限制长度
        if len(summary) > max_length:
            summary = summary[:max_length] + "..."
        
        return {
            "summary": summary,
            "summary_type": summary_type,
            "original_length": len(text_content),
            "summary_length": len(summary),
            "compression_ratio": len(summary) / len(text_content) if text_content else 0,
            "metadata": metadata
        }
    
    async def extract_keywords(self, file_path: Path, file_ext: str, parameters: Dict) -> Dict[str, Any]:
        """提取关键词"""
        text_content, _, _, _ = await self.extract_text(file_path, file_ext, {"extract_type": "text"})
        
        if not text_content:
            return {"error": "无法提取文档内容"}
        
        max_keywords = parameters.get("max_keywords", 20)
        min_freq = parameters.get("min_frequency", 2)
        
        # 简单的关键词提取（基于词频）
        if parameters.get("language", "zh") == "zh":
            # 中文分词
            words = jieba.lcut(text_content)
            # 过滤停用词和短词
            words = [word for word in words if len(word) > 1 and word not in self.get_chinese_stopwords()]
        else:
            # 英文分词
            words = text_content.lower().split()
            words = [word for word in words if len(word) > 2 and word.isalpha()]
        
        # 统计词频
        word_freq = Counter(words)
        
        # 过滤低频词
        filtered_words = {word: freq for word, freq in word_freq.items() if freq >= min_freq}
        
        # 排序并取前N个
        top_keywords = sorted(filtered_words.items(), key=lambda x: x[1], reverse=True)[:max_keywords]
        
        return {
            "keywords": [{"word": word, "frequency": freq} for word, freq in top_keywords],
            "total_words": len(words),
            "unique_words": len(word_freq),
            "extraction_method": "frequency_based"
        }
    
    def get_chinese_stopwords(self):
        """获取中文停用词列表"""
        return {'的', '了', '在', '是', '我', '有', '和', '就', '不', '人', '都', '一', '一个', '上', '也', '很', '到', '说', '要', '去', '你', '会', '着', '没有', '看', '好', '自己', '这'}
    
    async def classify_document(self, file_path: Path, file_ext: str, parameters: Dict) -> Dict[str, Any]:
        """文档分类"""
        text_content, _, _, metadata = await self.extract_text(file_path, file_ext, {"extract_type": "text"})
        
        if not text_content:
            return {"error": "无法提取文档内容"}
        
        # 简单的基于关键词的分类
        categories = {
            "技术文档": ["技术", "系统", "开发", "设计", "架构", "接口", "API", "数据库"],
            "测试报告": ["测试", "报告", "用例", "结果", "验证", "功能", "性能"],
            "用户手册": ["用户", "手册", "操作", "使用", "指南", "步骤", "说明"],
            "项目文档": ["项目", "计划", "进度", "需求", "方案", "总结", "汇报"],
            "失效分析": ["失效", "故障", "缺陷", "问题", "原因", "分析", "检测"]
        }
        
        scores = {}
        for category, keywords in categories.items():
            score = sum(1 for keyword in keywords if keyword in text_content)
            scores[category] = score
        
        # 找到最高分的分类
        best_category = max(scores.items(), key=lambda x: x[1])
        
        return {
            "predicted_category": best_category[0] if best_category[1] > 0 else "未分类",
            "confidence": best_category[1] / len(text_content.split()) if text_content else 0,
            "all_scores": scores,
            "classification_method": "keyword_based"
        }
    
    async def extract_entities(self, file_path: Path, file_ext: str, parameters: Dict) -> Dict[str, Any]:
        """实体提取"""
        text_content, _, _, _ = await self.extract_text(file_path, file_ext, {"extract_type": "text"})
        
        if not text_content:
            return {"error": "无法提取文档内容"}
        
        # 简单的实体提取（基于正则表达式）
        entities = {
            "dates": [],
            "numbers": [],
            "emails": [],
            "urls": [],
            "phone_numbers": []
        }
        
        # 日期
        date_pattern = r'\d{4}[-/年]\d{1,2}[-/月]\d{1,2}[日]?'
        entities["dates"] = list(set(re.findall(date_pattern, text_content)))
        
        # 数字
        number_pattern = r'\d+\.?\d*'
        entities["numbers"] = list(set(re.findall(number_pattern, text_content)))
        
        # 邮箱
        email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        entities["emails"] = list(set(re.findall(email_pattern, text_content)))
        
        # URL
        url_pattern = r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+'
        entities["urls"] = list(set(re.findall(url_pattern, text_content)))
        
        # 电话号码
        phone_pattern = r'1[3-9]\d{9}'
        entities["phone_numbers"] = list(set(re.findall(phone_pattern, text_content)))
        
        return {
            "entities": entities,
            "total_entities": sum(len(v) for v in entities.values()),
            "extraction_method": "regex_based"
        }
    
    async def calculate_similarity(self, file_path: Path, file_ext: str, parameters: Dict) -> Dict[str, Any]:
        """计算文档相似度"""
        text_content, _, _, _ = await self.extract_text(file_path, file_ext, {"extract_type": "text"})
        
        if not text_content:
            return {"error": "无法提取文档内容"}
        
        compare_text = parameters.get("compare_text", "")
        if not compare_text:
            return {"error": "缺少对比文本"}
        
        # 简单的相似度计算（基于词汇重叠）
        words1 = set(text_content.lower().split())
        words2 = set(compare_text.lower().split())
        
        intersection = words1.intersection(words2)
        union = words1.union(words2)
        
        jaccard_similarity = len(intersection) / len(union) if union else 0
        
        return {
            "similarity_score": jaccard_similarity,
            "similarity_method": "jaccard",
            "common_words": list(intersection)[:20],  # 限制返回数量
            "unique_words_doc1": list(words1 - words2)[:20],
            "unique_words_doc2": list(words2 - words1)[:20]
        }

# 全局分析器实例
analyzer = DocumentAnalyzer()

# API端点
@app.get("/")
async def root():
    return {"message": "文档分析AI服务运行中", "version": "1.0.0"}

@app.get("/health")
async def health_check():
    return {"status": "healthy", "timestamp": datetime.now().isoformat()}

@app.post("/analyze", response_model=DocumentAnalysisResult)
async def analyze_document(
    request: DocumentAnalysisRequest,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...)
):
    """分析文档"""
    start_time = datetime.now()
    
    try:
        # 验证文件类型
        if not file.filename:
            raise HTTPException(status_code=400, detail="文件名不能为空")
        
        file_ext = Path(file.filename).suffix.lower()
        if file_ext not in analyzer.supported_formats:
            raise HTTPException(status_code=400, detail=f"不支持的文件格式: {file_ext}")
        
        # 读取文件数据
        file_data = await file.read()
        
        # 更新任务状态
        await redis_client.set(f"task:{request.task_id}:status", "running")
        
        # 执行分析
        results = await analyzer.analyze_document(file_data, file.filename, request.analysis_type, request.parameters)
        
        # 计算执行时间
        execution_time = (datetime.now() - start_time).total_seconds()
        
        # 构建结果
        result = DocumentAnalysisResult(
            task_id=request.task_id,
            status="completed",
            results=results,
            execution_time=execution_time,
            timestamp=datetime.now().isoformat()
        )
        
        # 更新Redis中的结果
        await redis_client.setex(
            f"task:{request.task_id}:result",
            3600,  # 1小时过期
            json.dumps(result.dict())
        )
        
        return result
        
    except Exception as e:
        logger.error(f"文档分析失败: {str(e)}")
        await redis_client.set(f"task:{request.task_id}:status", "failed")
        await redis_client.set(f"task:{request.task_id}:error", str(e))
        
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/task/{task_id}/status")
async def get_task_status(task_id: str):
    """获取任务状态"""
    status = await redis_client.get(f"task:{task_id}:status")
    if not status:
        raise HTTPException(status_code=404, detail="任务不存在")
    
    result = {"task_id": task_id, "status": status}
    
    # 如果任务完成，返回结果
    if status == "completed":
        result_data = await redis_client.get(f"task:{task_id}:result")
        if result_data:
            result["result"] = json.loads(result_data)
    
    # 如果任务失败，返回错误信息
    elif status == "failed":
        error = await redis_client.get(f"task:{task_id}:error")
        if error:
            result["error"] = error
    
    return result

@app.get("/capabilities")
async def get_capabilities():
    """获取服务能力"""
    return {
        "supported_analysis_types": [
            "text_extraction",
            "document_summary",
            "keyword_extraction",
            "document_classification",
            "entity_extraction",
            "document_similarity"
        ],
        "supported_formats": analyzer.supported_formats,
        "max_file_size": "100MB",
        "ocr_enabled": ocr_reader is not None,
        "languages_supported": ["zh", "en", "auto"]
    }

if __name__ == "__main__":
    uvicorn.run(
        "main:app",
        host="0.0.0.0",
        port=8002,
        reload=True,
        log_level="info"
    )